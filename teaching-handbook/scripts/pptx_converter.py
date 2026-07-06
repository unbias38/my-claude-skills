# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx",
#     "beautifulsoup4",
# ]
# ///
import os
import sys
import argparse
import base64
import tempfile
import html as html_lib
from pptx import Presentation
from pptx.util import Emu

try:
    import style_injector
    from _polish import SIDEBAR_SEARCH_SCRIPT, add_lazy_loading
except ImportError:
    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
    import style_injector
    from _polish import SIDEBAR_SEARCH_SCRIPT, add_lazy_loading


def _image_to_data_uri(image_part):
    content_type = (image_part.content_type or "").lower()
    if any(fmt in content_type for fmt in ("emf", "wmf", "tiff")):
        print(f"  skipped unsupported image format: {content_type}")
        return None
    ext = (content_type.split("/")[-1] or "png").lower()
    if ext == "jpeg":
        ext = "jpeg"
    b64 = base64.b64encode(image_part.blob).decode("ascii")
    return f"data:image/{ext};base64,{b64}"


def _escape(text):
    return html_lib.escape(text or "", quote=False)


def _slide_title(slide, index):
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip()
        if t:
            return t
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                first_line = t.splitlines()[0].strip()
                if first_line:
                    return first_line
    return f"Slide {index}"


def _collect_body_paragraphs(slide, title_text):
    """Return (list_items, long_paragraphs) extracted from non-title text frames."""
    list_items = []
    long_paragraphs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape is slide.shapes.title:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                continue
            if text == title_text:
                continue
            # Skip lone numbers / single chars (usually slide page numbers)
            if len(text) <= 2 and (text.isdigit() or text.isalpha()):
                continue
            if len(text) <= 120:
                list_items.append(text)
            else:
                long_paragraphs.append(text)
    return list_items, long_paragraphs


def _collect_images(slide):
    data_uris = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                uri = _image_to_data_uri(shape.image)
                if uri is not None:
                    data_uris.append(uri)
            except Exception as e:
                print(f"  image extract failed: {e}")
    return data_uris


def _collect_notes(slide):
    if not slide.has_notes_slide:
        return ""
    nt = slide.notes_slide.notes_text_frame
    if nt is None:
        return ""
    return nt.text.strip()


def _is_chapter_slide(bullets, long_paras, images, layout_name):
    """A chapter divider has only a title (no body), or its layout signals so.

    Used to promote the slide's <h2> to <h1> so style_injector's sidebar
    renders it as a top-level group divider.
    """
    layout_lc = (layout_name or "").lower()
    if any(kw in layout_lc for kw in ("section header", "section divider", "title slide", "chapter")):
        return True
    if not bullets and not long_paras and len(images) <= 1:
        return True
    return False


def _render_slide_html(index, title, images, bullets, long_paras, notes, is_chapter=False):
    parts = [f'<section class="pptx-slide{" pptx-chapter-section" if is_chapter else ""}">']
    if is_chapter:
        # H1 is always picked up by style_injector's sidebar JS regardless of regex.
        parts.append(f'<h1 class="pptx-chapter">{_escape(title)}</h1>')
    else:
        # Prefix index so style_injector's sidebar regex (^\d+\s) picks it up.
        labeled = f"{index:02d} - {title}"
        parts.append(f'<h2>{_escape(labeled)}</h2>')

    if images:
        parts.append('<div class="pptx-thumbs">')
        for uri in images:
            parts.append(f'<img src="{uri}" alt="" loading="lazy">')
        parts.append('</div>')

    if bullets:
        parts.append("<ul>")
        for item in bullets:
            parts.append(f"<li>{_escape(item)}</li>")
        parts.append("</ul>")

    for para in long_paras:
        parts.append(f"<p>{_escape(para)}</p>")

    if notes:
        parts.append('<div class="pptx-notes">')
        parts.append("<p><strong>講者備註</strong></p>")
        for line in notes.splitlines():
            line = line.strip()
            if line:
                parts.append(f"<p>{_escape(line)}</p>")
        parts.append("</div>")

    parts.append("</section>")
    return "\n".join(parts)


def convert_pptx_to_html(pptx_path, output_filename, page_title="教學手冊",
                        sidebar_title="投影片目錄", include_notes=True):
    print(f"Converting {pptx_path} to {output_filename}...")

    if not os.path.exists(pptx_path):
        print(f"Error: Input file '{pptx_path}' not found.")
        sys.exit(1)

    prs = Presentation(pptx_path)

    sections = []
    chapter_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide, i)
        bullets, long_paras = _collect_body_paragraphs(slide, title)
        images = _collect_images(slide)
        notes = _collect_notes(slide) if include_notes else ""
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:
            pass
        is_chapter = _is_chapter_slide(bullets, long_paras, images, layout_name)
        if is_chapter:
            chapter_count += 1
        sections.append(_render_slide_html(i, title, images, bullets, long_paras, notes, is_chapter=is_chapter))
        marker = "[CH]" if is_chapter else "    "
        print(f"  {marker} slide {i}: {title[:40]}  (bullets={len(bullets)}, images={len(images)}, notes={'y' if notes else 'n'})")
    print(f"  Detected {chapter_count} chapter slide(s) -> promoted to <h1>")

    body = "\n".join(sections)

    raw_html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{_escape(page_title)}</title>
    <style>
        img {{ max-width: 100%; height: auto; }}
        .pptx-slide {{ margin-bottom: 2.5rem; padding-bottom: 1rem; border-bottom: 1px dashed #ccc; }}
        .pptx-chapter-section {{ margin-top: 2.5rem; padding-top: 0.5rem; border-bottom: none; }}
        h1.pptx-chapter {{ margin: 1rem 0 0.5rem; padding: 0.4rem 0 0.6rem; border-top: 3px solid #888; border-bottom: 1px solid #ddd; font-size: 1.7em; color: #2a2a2a; }}
        .pptx-thumbs {{ display: flex; flex-wrap: wrap; gap: 8px; margin: 0.5rem 0 1rem; }}
        .pptx-thumbs img {{ width: 80px; height: 80px; object-fit: contain; border: 1px solid #eee; border-radius: 4px; background: #fafafa; padding: 4px; }}
        .pptx-notes {{ margin-top: 0.75rem; padding: 0.75rem 1rem; background: #f7f7f5; border-left: 3px solid #bbb; font-size: 0.95em; }}
    </style>
</head>
<body>
    <div class="WordSection1">
    {body}
    </div>
{SIDEBAR_SEARCH_SCRIPT}
</body>
</html>
"""

    with tempfile.NamedTemporaryFile(
        mode="w", encoding="utf-8", suffix=".html", delete=False
    ) as f:
        temp_file = f.name
        f.write(raw_html)
    print(f"Created temporary file: {temp_file}")

    try:
        print("Injecting styles and navigation...")
        style_injector.inject_styles_and_nav(
            temp_file,
            output_filename,
            layout_mode="sidebar",
            page_title=page_title,
            sidebar_title=sidebar_title,
        )
    finally:
        if os.path.exists(temp_file):
            os.remove(temp_file)
            print("Cleaned up temp file.")

    if not os.path.exists(output_filename):
        print("Error: style injection failed, no output produced.")
        sys.exit(1)

    print(f"Done! Created {output_filename}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PPTX to High-Fidelity HTML with Sidebar")
    parser.add_argument("input", help="Input .pptx file path")
    parser.add_argument("output", nargs="?", help="Output .html file path (optional)")
    parser.add_argument("--title", default="教學手冊", help="Page title")
    parser.add_argument("--sidebar-title", default="投影片目錄", help="Sidebar navigation title")
    parser.add_argument("--no-notes", action="store_true", help="Skip speaker notes")

    args = parser.parse_args()

    output_path = args.output
    if not output_path:
        base_name = os.path.splitext(args.input)[0]
        output_path = f"{base_name}.html"

    convert_pptx_to_html(
        args.input,
        output_path,
        page_title=args.title,
        sidebar_title=args.sidebar_title,
        include_notes=not args.no_notes,
    )
